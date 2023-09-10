'''
Purpose : Create a presentation library which help create or modify a presentation
Author : Parth Gandhi
'''
import pandas as pd
import numpy as np
import os, time, re
from pptx import Presentation
#from pptx.util import Inches, Pt
from pptx.util import *

def create_new_ppt(path = str, filename = str):
    '''
    Create a new presentation
    '''
    # Create a new presentation
    presentation = Presentation()
    
    # Save the presentation
    output_path = os.path.join(path,filename)
    presentation.save(output_path)
    
    print('\nNew presentation created : %s'%(output_path))
    
    return presentation

def insert_slide(ppt = object, layout = int, position = int):
    '''
    Function to add slide of choice
    
    Title Slide (0):
    Typically used for the title slide at the beginning of a presentation.

    Title and Content (1):
    A layout with a title area and a content area, often used for regular content slides.

    Title Slide with Content (5):
    A layout with a title area and a larger content area, useful for presenting content with a title.

    Title, Content, and 2 Content (5):
    A layout with a title area and multiple content areas, suitable for more complex content slides.

    Title, Content, and Chart (6):
    A layout with a title area, content area, and a chart placeholder.

    Title, Content, and Table (7):
    A layout with a title area, content area, and a table placeholder.

    Title, Content, and SmartArt (8):
    A layout with a title area, content area, and a SmartArt placeholder.

    Title, Content, and Media (9):
    A layout with a title area, content area, and a media (e.g., video) placeholder.

    Content Slide (10):
    A layout with just a content area, useful for slides with large content.

    Blank (11):
    A completely blank layout with no predefined placeholders.
    
    Two Content (13):
    Features two content text placeholders without a title.

    Comparison (16):
    Offers two content text placeholders with a title, often used for comparisons.

    Title Only (17):
    Contains only a title text placeholder without any content.

    Blank Slide (18):
    A completely blank layout with no placeholders.

    Content with Caption (19):
    Contains a content text placeholder and a caption text placeholder.

    Picture with Caption (20):
    Provides a picture placeholder and a caption text placeholder.

    Section Header (21):
    Used as a section header slide.

    Two Content with Caption (22):
    Features two content text placeholders and two caption text placeholders.

    Picture with Caption (23):
    Offers a picture placeholder and two caption text placeholders.

    '''
    # Load the existing PowerPoint presentation
    presentation = ppt
    '''
    # Choose a slide layout for the new slide (e.g., index 1 for Title Slide layout)
    blank_slide_layout   = presentation.slide_layouts[layout]
    new_slide            = presentation.slides.add_slide(blank_slide_layout)
    presentation.slides._sldIdLst.insert(position, new_slide._element)
    '''
    # Choose a slide layout for the new slide
    blank_slide_layout = presentation.slide_layouts[layout]
    new_slide = presentation.slides.add_slide(blank_slide_layout)

    # Move the new slide to the desired position
    slides = list(presentation.slides)
    slides.insert(position, new_slide)
    for idx, slide in enumerate(slides):
        slide._element.rId = idx + 256

    return new_slide
    
    
    # Save the modified presentation
    print("\nAdded slide\n")
    
    return new_slide

def add_title(slide = object, title_txt = str, subtitle_txt = None):
    '''
    Function to add Title or Subtitle
    '''
    # Create/Assign space for title/subtitle
    title = slide.shapes.title
    
    # Add tile/subtitle text
    title.text = title_txt
    if subtitle_txt != None:
        subtitle = slide.placeholders[1]  # Assuming there's a subtitle placeholder
        subtitle.text = subtitle_txt
    
    print("\nAdded tile/subtitle\n")

def add_image(ppt = object, layout = int, position = int, img_path = str):
    '''
    Function to add single image on the presentation foil
    '''
    
   # Load the existing PowerPoint presentation
    presentation = ppt

    # Add a slide with the chosen layout
    slide =  insert_slide(ppt = presentation, layout = layout, position = position)
    
    # Paths to the images you want to add
    image_path = img_path
    
    # Add the image to the center of the slide
    left     = Inches(0.5)                              # Adjust the left position as needed
    top      = Inches(1.5)                              # Adjust the top position as needed
    width    = presentation.slide_width - Inches(1)     # Auto fit width to slide
    height   = presentation.slide_height - Inches(3)    # Auto fit height to slide
    
    # Add images to the slide
    slide.shapes.add_picture(image_path, left, top, width, height)
    
    # Exit function
    print('\nAdded Image to the slide\n')

def add_multiple_image(ppt = object, layout = int, position = int, img_path = list, img_dim = list, img_pos = list):
    '''
    Function to add multiple images on the presentation foil
    '''
    # Load the existing PowerPoint presentation
    presentation = ppt
    
    # Add a slide with the chosen layout
    slide =  insert_slide(ppt = presentation, layout = layout, position = position)
    
    # Paths to the images you want to add
    image_paths = img_path
    
    # Coordinates and dimensions for each image (in Inches)
    #image_positions  = [(0.5, 1), (3, 1), (5.5, 1)]
    #image_dimensions = [(2, 2), (2, 2), (2, 2)]
    image_positions  = img_pos
    image_dimensions = img_dim
    
    # Define the position, size and add images
    for i, image_path in enumerate(image_paths):
        
        # Define the position and size of images (in Inches)
        h, v = image_positions[i]
        w, ht = image_dimensions[i]
        print(i, image_positions[i], h,v)
        
        if h >= 0:
            horizontal     = Inches(h)
        else :
            horizontal     = presentation.slide_width - Inches(np.abs(h))
    
        if v >= 0:
            vertical     = Inches(v)
        else :
            vertical     = presentation.slide_height - Inches(np.abs(v))
        
        width    = Inches(w)
        height   = Inches(ht)
        
        # Add images to the slide
        slide.shapes.add_picture(image_path, horizontal, vertical, width, height)
        print(image_path, horizontal, vertical, width, height)

    # Save the presentation
    #presentation.save("slide_with_images.pptx")
    
    # Exit function
    print('\nAdded Images to the slide\n')

def add_text_box(ppt = object, slide = object, txt = str, text_pos = list):
    '''
    Function to add Textbox and text
    '''
   # Load the existing PowerPoint presentation
    presentation = ppt
    
    # Define the position and size of the text box (in Inches)
    if text_pos[0] >= 0:
        horizontal     = Inches(text_pos[0])
    else :
        horizontal     = presentation.slide_width - Inches(np.abs(text_pos[0]))
    
    if text_pos[1] >= 0:
        vertical     = Inches(text_pos[1])
    else :
        vertical     = presentation.slide_height - Inches(np.abs(text_pos[1]))
    
    width    = Inches(text_pos[2])
    height   = Inches(text_pos[3])
    
    
    # Add a text box to the slide
    text_box    = slide.shapes.add_textbox(horizontal, vertical, width, height)
    text_frame  = text_box.text_frame

    # Add content to the text box
    text_frame.text = txt

    # Customize the font of the text box
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            font.name = "Calibri"
            font.size = Pt(10)
    
    print("\nAdded text box and text\n")

def main(ppt_path=r'', layout = int, position = int):
    
    # Load the existing PowerPoint presentation
    existing_pptx = ppt_path
    presentation = Presentation(existing_pptx)
    
    # Add slide to existing presentation
    new_slide = insert_slide(ppt = presentation, layout = int, position = int)
    
    # Add title and subtitle:
    